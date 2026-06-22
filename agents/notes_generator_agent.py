"""notes_generator_agent.py — LLM-powered notes and document generation plugin.

Unlike tool-only agents, this agent **uses LLMs** (via Ollama / OpenAI-compatible
APIs) to generate educational, technical, academic, professional, and
documentation content. The Brain Agent decides *when* notes are needed and
*what* topic / structure to request; this agent handles content creation and
document formatting.

Dependencies (install once)::

    pip install ollama python-docx reportlab python-pptx PyPDF2 python-dotenv

Environment variables (read from ``.env`` via python-dotenv)::

    NOTES_MODEL         — Ollama model name  (default: ``gpt-oss:120b``)
    OLLAMA_BASE_URL     — Ollama server URL  (default: ``http://localhost:11434``)
    DEFAULT_OUTPUT_DIR  — output folder      (default: ``generated_notes``)

CLI usage::

    python notes_generator_agent.py request.json
    python notes_generator_agent.py '{"action":"generate_notes","parameters":{"topic":"Docker"}}'
"""

from __future__ import annotations

import json
import math
import os
import re
import sys
import textwrap
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Optional

from dotenv import load_dotenv

__all__ = ["execute", "PLUGIN_INFO"]


# --------------------------------------------------------------------------- #
# Plugin metadata
# --------------------------------------------------------------------------- #
PLUGIN_INFO: dict[str, Any] = {
    "name": "notes_generator_agent",
    "agent_name": "NotesGeneratorAgent",
    "version": "1.0.0",
    "type": "content_generator",
    "input_format": "json",
    "output_format": "json",
    "entrypoint": "execute",
}


# --------------------------------------------------------------------------- #
# Configuration
# --------------------------------------------------------------------------- #
_BASE_DIR: Path = Path(os.path.dirname(os.path.abspath(__file__)))

load_dotenv()
load_dotenv(_BASE_DIR / ".env", override=False)

_MASTER_MODEL: str = os.getenv("MAYA_MODEL", "qwen3:8b")
NOTES_MODEL: str = os.getenv("NOTES_MODEL", _MASTER_MODEL)
OLLAMA_BASE_URL: str = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
DEFAULT_OUTPUT_DIR: str = os.getenv("DEFAULT_OUTPUT_DIR", "generated_notes")

WORDS_PER_PAGE: int = 400
DEFAULT_PAGES: int = 10
MAX_PAGES: int = 200
DEFAULT_DIFFICULTY: str = "intermediate"
DEFAULT_AUDIENCE: str = "professional"
DEFAULT_FORMAT: str = "md"

VALID_DIFFICULTIES: frozenset[str] = frozenset({
    "beginner", "intermediate", "advanced", "expert",
})
VALID_AUDIENCES: frozenset[str] = frozenset({
    "school", "college", "engineering", "professional", "research", "executive",
})
VALID_FORMATS: frozenset[str] = frozenset({
    "docx", "pdf", "pptx", "md", "txt",
})
TEXT_EXTENSIONS: frozenset[str] = frozenset({
    ".txt", ".md", ".py", ".js", ".ts", ".html", ".css", ".xml",
    ".yaml", ".yml", ".json", ".csv",
})


# --------------------------------------------------------------------------- #
# Exceptions
# --------------------------------------------------------------------------- #
class NotesGeneratorError(Exception):
    """Raised for any handled error whose message is safe to return."""


# --------------------------------------------------------------------------- #
# LLM interaction
# --------------------------------------------------------------------------- #
def _llm_generate(prompt: str, *, model: Optional[str] = None) -> str:
    """Send *prompt* to the Ollama chat API and return the assistant reply."""
    import ollama as _ollama

    target_model = model or NOTES_MODEL
    try:
        client = _ollama.Client(host=OLLAMA_BASE_URL)
        response = client.chat(
            model=target_model,
            messages=[{"role": "user", "content": prompt}],
            options={"num_ctx": 16384},
        )
        return response.message.content or ""
    except Exception as exc:
        raise NotesGeneratorError(f"LLM generation failed ({target_model}): {exc}")


def _llm_generate_long(prompt: str, word_target: int, *, model: Optional[str] = None) -> str:
    """Generate content, requesting continuations until *word_target* is met."""
    accumulated = _llm_generate(prompt, model=model)
    iteration = 0
    max_iterations = 15

    while len(accumulated.split()) < word_target * 0.85 and iteration < max_iterations:
        iteration += 1
        cont_prompt = (
            f"Continue writing from where you left off. "
            f"You have written {len(accumulated.split())} words so far. "
            f"Target is {word_target} words. Continue the content seamlessly:\n\n"
            f"...{accumulated[-1500:]}"
        )
        chunk = _llm_generate(cont_prompt, model=model)
        if not chunk.strip():
            break
        accumulated += "\n\n" + chunk

    return accumulated


# --------------------------------------------------------------------------- #
# Prompt builders
# --------------------------------------------------------------------------- #
def _build_notes_prompt(
    topic: str,
    word_count: int,
    difficulty: str,
    audience: str,
    structure: Optional[list[str]] = None,
    academic_mode: bool = False,
    extra_instructions: str = "",
) -> str:
    """Build the system prompt for notes generation."""
    structure_block = ""
    if structure:
        numbered = "\n".join(f"  {i}. {s}" for i, s in enumerate(structure, 1))
        structure_block = f"\nFollow this exact structure:\n{numbered}\n"

    academic_block = ""
    if academic_mode:
        academic_block = textwrap.dedent("""
        Include academic elements:
        - Module-wise / Unit-wise organisation
        - Important definitions
        - Key formulas or frameworks
        - Advantages and disadvantages
        - Real-world applications
        - Important exam questions with brief answers
        - Viva / interview questions
        - Key takeaways per section
        """)

    return textwrap.dedent(f"""\
        Write comprehensive, detailed notes on the topic: "{topic}"

        Requirements:
        - Target length: approximately {word_count} words
        - Difficulty level: {difficulty}
        - Target audience: {audience}
        - Use clear headings and subheadings (Markdown format)
        - Include detailed explanations with real-world examples
        - Include tables where appropriate (Markdown tables)
        - Include diagrams descriptions where helpful
        - Include important notes and warnings where relevant
        - Include key takeaways at the end of each major section
        {structure_block}{academic_block}
        {extra_instructions}
        Write in a clear, professional, educational tone.
        Do NOT include meta-commentary about the writing process.
        Start the content directly.
    """).strip()


def _build_exam_prompt(topic: str, word_count: int, difficulty: str, audience: str) -> str:
    return _build_notes_prompt(
        topic, word_count, difficulty, audience,
        academic_mode=True,
        extra_instructions=(
            "Focus on exam preparation. Include:\n"
            "- Likely exam questions (short answer, long answer, MCQ-style)\n"
            "- Model answers for each question\n"
            "- Mnemonics or memory aids where applicable\n"
            "- Common mistakes students make\n"
            "- Mark-scheme style point breakdowns"
        ),
    )


def _build_short_prompt(topic: str, word_count: int, difficulty: str, audience: str) -> str:
    return textwrap.dedent(f"""\
        Write concise, condensed notes on: "{topic}"

        Requirements:
        - Target length: approximately {word_count} words
        - Difficulty: {difficulty}, Audience: {audience}
        - Use bullet points and short paragraphs
        - Include only the most essential information
        - Use Markdown headings for organisation
        - Include a quick-reference summary table
        - Include key definitions in bold
        - End with a one-page cheat-sheet section

        Be concise but accurate. Start the content directly.
    """).strip()


def _build_presentation_prompt(topic: str, slide_count: int, difficulty: str, audience: str) -> str:
    return textwrap.dedent(f"""\
        Create presentation content for: "{topic}"

        Requirements:
        - Create exactly {slide_count} slides
        - Difficulty: {difficulty}, Audience: {audience}
        - For each slide provide:
          SLIDE <number>: <Title>
          CONTENT:
          - Bullet point 1
          - Bullet point 2
          - Bullet point 3
          SPEAKER_NOTES: <paragraph of speaker notes>
        - Include a title slide and a summary/conclusion slide
        - Keep bullet points concise (max 8 words each)
        - Include 3-6 bullet points per slide
        - Speaker notes should elaborate on the bullet points

        Start directly with SLIDE 1. Do not add preamble.
    """).strip()


def _build_report_prompt(topic: str, word_count: int, difficulty: str, audience: str) -> str:
    return textwrap.dedent(f"""\
        Write a professional report on: "{topic}"

        Requirements:
        - Target length: approximately {word_count} words
        - Difficulty: {difficulty}, Audience: {audience}
        - Include: Executive Summary, Introduction, detailed body sections,
          Findings/Analysis sections, Conclusions, Recommendations
        - Use formal, professional language
        - Include data references and statistics where relevant
        - Use Markdown headings and formatting
        - Include tables for comparative data

        Start the report directly. No meta-commentary.
    """).strip()


# --------------------------------------------------------------------------- #
# Source file reading
# --------------------------------------------------------------------------- #
def _read_source_file(path_str: str) -> str:
    """Read and extract text from a source file for note generation."""
    path = Path(path_str).resolve()
    if not path.exists():
        raise NotesGeneratorError(f"Source file not found: {path}")
    if not path.is_file():
        raise NotesGeneratorError(f"Path is not a file: {path}")

    ext = path.suffix.lower()

    if ext in TEXT_EXTENSIONS:
        try:
            return path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            return path.read_text(encoding="latin-1")

    if ext == ".pdf":
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            raise NotesGeneratorError("PyPDF2 not installed. Run: pip install PyPDF2")
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        if not pages:
            raise NotesGeneratorError("No text could be extracted from PDF")
        return "\n\n".join(pages)

    if ext == ".docx":
        try:
            from docx import Document
        except ImportError:
            raise NotesGeneratorError("python-docx not installed. Run: pip install python-docx")
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            raise NotesGeneratorError("python-pptx not installed. Run: pip install python-pptx")
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        if not texts:
            raise NotesGeneratorError("No text could be extracted from PPTX")
        return "\n".join(texts)

    raise NotesGeneratorError(f"Unsupported source file type: {ext}")


# --------------------------------------------------------------------------- #
# Document writers
# --------------------------------------------------------------------------- #
def _output_dir() -> Path:
    d = _BASE_DIR / DEFAULT_OUTPUT_DIR
    d.mkdir(parents=True, exist_ok=True)
    return d


def _safe_filename(topic: str) -> str:
    """Convert a topic into a filesystem-safe slug."""
    slug = re.sub(r"[^\w\s-]", "", topic.lower())
    slug = re.sub(r"[\s_-]+", "_", slug).strip("_")
    return slug[:80] or "notes"


def _write_md(content: str, topic: str) -> Path:
    name = _safe_filename(topic) + ".md"
    path = _output_dir() / name
    path.write_text(content, encoding="utf-8")
    return path


def _write_txt(content: str, topic: str) -> Path:
    plain = content
    name = _safe_filename(topic) + ".txt"
    path = _output_dir() / name
    path.write_text(plain, encoding="utf-8")
    return path


def _write_docx(content: str, topic: str) -> Path:
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        raise NotesGeneratorError("python-docx not installed. Run: pip install python-docx")

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    lines = content.split("\n")
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("#### "):
            doc.add_heading(stripped[5:], level=4)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        elif re.match(r"^\d+\.\s", stripped):
            text = re.sub(r"^\d+\.\s*", "", stripped)
            doc.add_paragraph(text, style="List Number")
        elif stripped.startswith("|") and stripped.endswith("|"):
            doc.add_paragraph(stripped)
        else:
            doc.add_paragraph(stripped)

    name = _safe_filename(topic) + ".docx"
    path = _output_dir() / name
    doc.save(str(path))
    return path


def _write_pdf(content: str, topic: str) -> Path:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            Paragraph,
            SimpleDocTemplate,
            Spacer,
        )
    except ImportError:
        raise NotesGeneratorError("reportlab not installed. Run: pip install reportlab")

    name = _safe_filename(topic) + ".pdf"
    path = _output_dir() / name
    doc = SimpleDocTemplate(
        str(path), pagesize=A4,
        leftMargin=inch, rightMargin=inch,
        topMargin=inch, bottomMargin=inch,
    )

    styles = getSampleStyleSheet()
    body_style = styles["BodyText"]
    body_style.fontSize = 11
    body_style.leading = 15

    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=18, spaceAfter=12)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=15, spaceAfter=10)
    h3 = ParagraphStyle("H3", parent=styles["Heading3"], fontSize=13, spaceAfter=8)
    bullet_style = ParagraphStyle("Bullet", parent=body_style, leftIndent=20, bulletIndent=10)

    story: list[Any] = []
    lines = content.split("\n")

    def _esc(text: str) -> str:
        return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    for line in lines:
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 6))
            continue
        if stripped.startswith("# "):
            story.append(Paragraph(_esc(stripped[2:]), h1))
        elif stripped.startswith("## "):
            story.append(Paragraph(_esc(stripped[3:]), h2))
        elif stripped.startswith("### ") or stripped.startswith("#### "):
            text = re.sub(r"^#{3,4}\s*", "", stripped)
            story.append(Paragraph(_esc(text), h3))
        elif stripped.startswith("- ") or stripped.startswith("* "):
            story.append(Paragraph(f"• {_esc(stripped[2:])}", bullet_style))
        elif re.match(r"^\d+\.\s", stripped):
            story.append(Paragraph(_esc(stripped), body_style))
        else:
            story.append(Paragraph(_esc(stripped), body_style))

    doc.build(story)
    return path


def _parse_slides(content: str) -> list[dict[str, Any]]:
    """Parse LLM output into slide dicts with title, bullets, speaker_notes.

    Handles two formats:
    - Explicit ``SLIDE N: Title`` format (preferred)
    - Markdown headings (``# Title`` / ``## Title``) as fallback
    """
    slides: list[dict[str, Any]] = []
    current: Optional[dict[str, Any]] = None

    section = "bullets"
    for line in content.split("\n"):
        stripped = line.strip()

        slide_match = re.match(r"^SLIDE\s+\d+\s*:\s*(.+)", stripped, re.IGNORECASE)
        if not slide_match:
            slide_match = re.match(r"^#{1,2}\s+(.+)", stripped)

        if slide_match:
            if current:
                slides.append(current)
            current = {"title": slide_match.group(1).strip(), "bullets": [], "speaker_notes": ""}
            section = "bullets"
            continue
        if current is None:
            continue
        if stripped.upper().startswith("CONTENT:"):
            section = "bullets"
            continue
        if stripped.upper().startswith("SPEAKER_NOTES:") or stripped.upper().startswith("SPEAKER NOTES:"):
            section = "notes"
            rest = re.sub(r"^SPEAKER[_ ]?NOTES:\s*", "", stripped, flags=re.IGNORECASE)
            if rest:
                current["speaker_notes"] = rest
            continue
        if re.match(r"^#{3,}\s+", stripped):
            bullet = re.sub(r"^#{3,}\s*", "", stripped)
            if bullet:
                current["bullets"].append(bullet)
            continue
        if section == "bullets":
            bullet = re.sub(r"^[-*•\d.]+[\s.)]*", "", stripped)
            if bullet:
                current["bullets"].append(bullet)
        elif section == "notes":
            if stripped:
                if current["speaker_notes"]:
                    current["speaker_notes"] += " " + stripped
                else:
                    current["speaker_notes"] = stripped

    if current:
        slides.append(current)
    return slides


def _write_pptx(content: str, topic: str) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise NotesGeneratorError("python-pptx not installed. Run: pip install python-pptx")

    slides_data = _parse_slides(content)
    if not slides_data:
        raise NotesGeneratorError("Could not parse presentation slides from generated content")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    for i, sd in enumerate(slides_data):
        layout = title_layout if i == 0 else content_layout
        slide = prs.slides.add_slide(layout)

        if slide.placeholders:
            slide.placeholders[0].text = sd["title"]

        if len(slide.placeholders) > 1 and sd["bullets"]:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for j, bullet in enumerate(sd["bullets"]):
                if j == 0:
                    tf.paragraphs[0].text = bullet
                    tf.paragraphs[0].font.size = Pt(18)
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(18)

        if sd["speaker_notes"]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = sd["speaker_notes"]

    name = _safe_filename(topic) + ".pptx"
    path = _output_dir() / name
    prs.save(str(path))
    return path


_WRITERS: dict[str, Callable[[str, str], Path]] = {
    "md": _write_md,
    "txt": _write_txt,
    "docx": _write_docx,
    "pdf": _write_pdf,
    "pptx": _write_pptx,
}


# --------------------------------------------------------------------------- #
# Param helpers
# --------------------------------------------------------------------------- #
def _get_topic(params: dict[str, Any]) -> str:
    topic = str(params.get("topic", "")).strip()
    if not topic:
        raise NotesGeneratorError("Missing topic parameter")
    return topic


def _get_pages(params: dict[str, Any]) -> int:
    raw = params.get("pages", DEFAULT_PAGES)
    try:
        pages = int(raw)
    except (TypeError, ValueError):
        raise NotesGeneratorError("pages must be an integer")
    return max(1, min(pages, MAX_PAGES))


def _get_word_count(params: dict[str, Any], pages: int) -> int:
    raw = params.get("word_count")
    if raw is not None:
        try:
            return max(100, int(raw))
        except (TypeError, ValueError):
            raise NotesGeneratorError("word_count must be an integer")
    return pages * WORDS_PER_PAGE


def _get_difficulty(params: dict[str, Any]) -> str:
    d = str(params.get("difficulty", DEFAULT_DIFFICULTY)).strip().lower()
    if d not in VALID_DIFFICULTIES:
        raise NotesGeneratorError(
            f"Invalid difficulty: {d}. Valid: {', '.join(sorted(VALID_DIFFICULTIES))}"
        )
    return d


def _get_audience(params: dict[str, Any]) -> str:
    a = str(params.get("audience", DEFAULT_AUDIENCE)).strip().lower()
    if a not in VALID_AUDIENCES:
        raise NotesGeneratorError(
            f"Invalid audience: {a}. Valid: {', '.join(sorted(VALID_AUDIENCES))}"
        )
    return a


def _get_format(params: dict[str, Any]) -> str:
    f = str(params.get("output_format", params.get("format", DEFAULT_FORMAT))).strip().lower()
    if f not in VALID_FORMATS:
        raise NotesGeneratorError(
            f"Invalid output_format: {f}. Valid: {', '.join(sorted(VALID_FORMATS))}"
        )
    return f


def _get_model(params: dict[str, Any]) -> Optional[str]:
    m = params.get("model")
    if m and isinstance(m, str) and m.strip():
        return m.strip()
    return None


def _write_output(content: str, topic: str, fmt: str) -> tuple[Path, int]:
    """Write content in the requested format and return (path, word_count)."""
    writer = _WRITERS.get(fmt)
    if not writer:
        raise NotesGeneratorError(f"No writer for format: {fmt}")
    path = writer(content, topic)
    wc = len(content.split())
    return path, wc


# --------------------------------------------------------------------------- #
# Action handlers
# --------------------------------------------------------------------------- #
def _action_generate_notes(params: dict[str, Any]) -> dict[str, Any]:
    topic = _get_topic(params)
    pages = _get_pages(params)
    word_count = _get_word_count(params, pages)
    difficulty = _get_difficulty(params)
    audience = _get_audience(params)
    fmt = _get_format(params)
    model = _get_model(params)
    structure = params.get("structure")
    academic = bool(params.get("academic_mode", False))

    prompt = _build_notes_prompt(topic, word_count, difficulty, audience, structure, academic)

    if fmt == "pptx":
        slide_count = max(5, pages * 2)
        prompt = _build_presentation_prompt(topic, slide_count, difficulty, audience)
        content = _llm_generate(prompt, model=model)
    else:
        content = _llm_generate_long(prompt, word_count, model=model)

    path, wc = _write_output(content, topic, fmt)
    return {
        "file_path": str(path),
        "pages_generated": max(1, math.ceil(wc / WORDS_PER_PAGE)),
        "word_count": wc,
        "topic": topic,
        "format": fmt,
    }


def _action_generate_exam_notes(params: dict[str, Any]) -> dict[str, Any]:
    topic = _get_topic(params)
    pages = _get_pages(params)
    word_count = _get_word_count(params, pages)
    difficulty = _get_difficulty(params)
    audience = _get_audience(params)
    fmt = _get_format(params)
    model = _get_model(params)

    prompt = _build_exam_prompt(topic, word_count, difficulty, audience)
    content = _llm_generate_long(prompt, word_count, model=model)
    path, wc = _write_output(content, topic, fmt)
    return {
        "file_path": str(path),
        "pages_generated": max(1, math.ceil(wc / WORDS_PER_PAGE)),
        "word_count": wc,
        "topic": topic,
        "format": fmt,
        "mode": "exam",
    }


def _action_generate_detailed_notes(params: dict[str, Any]) -> dict[str, Any]:
    topic = _get_topic(params)
    pages = _get_pages(params)
    word_count = _get_word_count(params, pages)
    difficulty = _get_difficulty(params)
    audience = _get_audience(params)
    fmt = _get_format(params)
    model = _get_model(params)

    prompt = _build_notes_prompt(
        topic, word_count, difficulty, audience,
        extra_instructions=(
            "Write extremely detailed, in-depth notes. "
            "Cover every subtopic thoroughly with multiple examples, "
            "case studies, and technical depth. "
            "Include diagrams descriptions, flowchart descriptions, "
            "and architecture explanations where relevant."
        ),
    )
    content = _llm_generate_long(prompt, word_count, model=model)
    path, wc = _write_output(content, topic, fmt)
    return {
        "file_path": str(path),
        "pages_generated": max(1, math.ceil(wc / WORDS_PER_PAGE)),
        "word_count": wc,
        "topic": topic,
        "format": fmt,
        "mode": "detailed",
    }


def _action_generate_short_notes(params: dict[str, Any]) -> dict[str, Any]:
    topic = _get_topic(params)
    pages = max(1, _get_pages(params) // 2)
    word_count = _get_word_count(params, pages)
    difficulty = _get_difficulty(params)
    audience = _get_audience(params)
    fmt = _get_format(params)
    model = _get_model(params)

    prompt = _build_short_prompt(topic, word_count, difficulty, audience)
    content = _llm_generate(prompt, model=model)
    path, wc = _write_output(content, topic, fmt)
    return {
        "file_path": str(path),
        "pages_generated": max(1, math.ceil(wc / WORDS_PER_PAGE)),
        "word_count": wc,
        "topic": topic,
        "format": fmt,
        "mode": "short",
    }


def _action_generate_presentation(params: dict[str, Any]) -> dict[str, Any]:
    topic = _get_topic(params)
    slides = int(params.get("slides", params.get("pages", 15)))
    slides = max(3, min(slides, 100))
    difficulty = _get_difficulty(params)
    audience = _get_audience(params)
    model = _get_model(params)

    prompt = _build_presentation_prompt(topic, slides, difficulty, audience)
    content = _llm_generate_long(prompt, slides * 80, model=model)
    path, wc = _write_output(content, topic, "pptx")
    return {
        "file_path": str(path),
        "slides_requested": slides,
        "word_count": wc,
        "topic": topic,
        "format": "pptx",
    }


def _action_generate_report(params: dict[str, Any]) -> dict[str, Any]:
    topic = _get_topic(params)
    pages = _get_pages(params)
    word_count = _get_word_count(params, pages)
    difficulty = _get_difficulty(params)
    audience = _get_audience(params)
    fmt = _get_format(params)
    model = _get_model(params)

    prompt = _build_report_prompt(topic, word_count, difficulty, audience)
    content = _llm_generate_long(prompt, word_count, model=model)
    path, wc = _write_output(content, topic, fmt)
    return {
        "file_path": str(path),
        "pages_generated": max(1, math.ceil(wc / WORDS_PER_PAGE)),
        "word_count": wc,
        "topic": topic,
        "format": fmt,
        "mode": "report",
    }


def _action_generate_from_structure(params: dict[str, Any]) -> dict[str, Any]:
    topic = _get_topic(params)
    structure = params.get("structure")
    if not structure or not isinstance(structure, list):
        raise NotesGeneratorError("Missing or invalid structure parameter (expected a list)")
    if not all(isinstance(s, str) and s.strip() for s in structure):
        raise NotesGeneratorError("All structure items must be non-empty strings")

    pages = _get_pages(params)
    word_count = _get_word_count(params, pages)
    difficulty = _get_difficulty(params)
    audience = _get_audience(params)
    fmt = _get_format(params)
    model = _get_model(params)

    prompt = _build_notes_prompt(topic, word_count, difficulty, audience, structure)
    content = _llm_generate_long(prompt, word_count, model=model)
    path, wc = _write_output(content, topic, fmt)
    return {
        "file_path": str(path),
        "pages_generated": max(1, math.ceil(wc / WORDS_PER_PAGE)),
        "word_count": wc,
        "topic": topic,
        "format": fmt,
        "structure_sections": len(structure),
    }


def _action_generate_from_prompt(params: dict[str, Any]) -> dict[str, Any]:
    prompt_text = str(params.get("prompt", "")).strip()
    if not prompt_text:
        raise NotesGeneratorError("Missing prompt parameter")

    fmt = _get_format(params)
    model = _get_model(params)
    topic = str(params.get("topic", "custom_notes")).strip()
    word_count = int(params.get("word_count", 2000))

    content = _llm_generate_long(prompt_text, word_count, model=model)
    path, wc = _write_output(content, topic, fmt)
    return {
        "file_path": str(path),
        "pages_generated": max(1, math.ceil(wc / WORDS_PER_PAGE)),
        "word_count": wc,
        "topic": topic,
        "format": fmt,
        "mode": "custom_prompt",
    }


def _action_generate_from_file(params: dict[str, Any]) -> dict[str, Any]:
    source_path = str(params.get("source_path", params.get("path", ""))).strip()
    if not source_path:
        raise NotesGeneratorError("Missing source_path parameter")

    source_content = _read_source_file(source_path)
    topic = str(params.get("topic", Path(source_path).stem)).strip()
    pages = _get_pages(params)
    word_count = _get_word_count(params, pages)
    difficulty = _get_difficulty(params)
    audience = _get_audience(params)
    fmt = _get_format(params)
    model = _get_model(params)

    max_source = 12000
    if len(source_content) > max_source:
        source_content = source_content[:max_source] + "\n\n[... content truncated ...]"

    prompt = textwrap.dedent(f"""\
        Based on the following source material, create comprehensive,
        well-structured notes on: "{topic}"

        Requirements:
        - Target length: approximately {word_count} words
        - Difficulty: {difficulty}, Audience: {audience}
        - Expand on the source material with additional explanations
        - Rewrite in a professional, clear, educational tone
        - Use Markdown headings and formatting
        - Include key takeaways per section
        - Add examples where the source material lacks them

        SOURCE MATERIAL:
        ---
        {source_content}
        ---

        Generate comprehensive notes based on this material.
        Start the content directly.
    """).strip()

    content = _llm_generate_long(prompt, word_count, model=model)
    path, wc = _write_output(content, topic, fmt)
    return {
        "file_path": str(path),
        "pages_generated": max(1, math.ceil(wc / WORDS_PER_PAGE)),
        "word_count": wc,
        "topic": topic,
        "format": fmt,
        "source_file": source_path,
        "mode": "from_file",
    }


# --------------------------------------------------------------------------- #
# Action registry
# --------------------------------------------------------------------------- #
_REQUIRED_PARAMS: dict[str, list[str]] = {
    "generate_notes": ["topic"],
    "generate_exam_notes": ["topic"],
    "generate_detailed_notes": ["topic"],
    "generate_short_notes": ["topic"],
    "generate_presentation": ["topic"],
    "generate_report": ["topic"],
    "generate_from_structure": ["topic", "structure"],
    "generate_from_prompt": ["prompt"],
    "generate_from_file": [],
}

_ACTIONS: dict[str, Callable[[dict[str, Any]], dict[str, Any]]] = {
    "generate_notes": _action_generate_notes,
    "generate_exam_notes": _action_generate_exam_notes,
    "generate_detailed_notes": _action_generate_detailed_notes,
    "generate_short_notes": _action_generate_short_notes,
    "generate_presentation": _action_generate_presentation,
    "generate_report": _action_generate_report,
    "generate_from_structure": _action_generate_from_structure,
    "generate_from_prompt": _action_generate_from_prompt,
    "generate_from_file": _action_generate_from_file,
}


# --------------------------------------------------------------------------- #
# Public entry point
# --------------------------------------------------------------------------- #
def execute(request_json: dict[str, Any]) -> dict[str, Any]:
    """Execute a notes-generation action.

    Args:
        request_json: A dict with ``"action"`` (str) and ``"parameters"`` (dict).

    Returns:
        A JSON-compatible dict with ``"status"`` (``"success"`` | ``"error"``),
        ``"action"``, and either ``"data"`` or ``"message"``.
    """
    if isinstance(request_json, str):
        request_json = request_json.lstrip("﻿")
        try:
            request_json = json.loads(request_json)
        except json.JSONDecodeError as exc:
            return {"status": "error", "action": "unknown", "message": f"Invalid JSON: {exc}"}

    if not isinstance(request_json, dict):
        return {"status": "error", "action": "unknown", "message": "Request must be a JSON object"}

    action = request_json.get("action", "")
    if not action or not isinstance(action, str):
        return {"status": "error", "action": "unknown", "message": "Missing or invalid action field"}

    if action not in _ACTIONS:
        return {
            "status": "error",
            "action": action,
            "message": f"Unknown action: {action}. Available: {', '.join(sorted(_ACTIONS))}",
        }

    parameters = request_json.get("parameters", {})
    if not isinstance(parameters, dict):
        return {"status": "error", "action": action, "message": "parameters must be a JSON object"}

    required = _REQUIRED_PARAMS.get(action, [])
    for key in required:
        value = parameters.get(key)
        if value is None or (isinstance(value, str) and not value.strip()):
            return {"status": "error", "action": action, "message": f"Missing required parameter: {key}"}

    try:
        data = _ACTIONS[action](parameters)
        return {"status": "success", "action": action, "data": data}
    except NotesGeneratorError as exc:
        return {"status": "error", "action": action, "message": str(exc)}
    except Exception as exc:
        return {"status": "error", "action": action, "message": f"Unexpected error: {exc}"}


# --------------------------------------------------------------------------- #
# CLI entry point
# --------------------------------------------------------------------------- #
def _cli_main() -> None:
    """Handle command-line invocation: file arg, inline JSON arg, or piped stdin."""
    raw: Optional[str] = None

    if len(sys.argv) > 1:
        arg = sys.argv[1]
        if os.path.isfile(arg):
            with open(arg, "rb") as fh:
                raw = fh.read().decode("utf-8-sig")
        else:
            raw = arg
    elif not sys.stdin.isatty():
        raw = sys.stdin.buffer.read().decode("utf-8-sig")
    else:
        print(
            "Usage:\n"
            '  python notes_generator_agent.py request.json\n'
            '  python notes_generator_agent.py \'{"action":"generate_notes",'
            '"parameters":{"topic":"Docker"}}\'\n'
            '  echo \'{"action":"generate_short_notes",'
            '"parameters":{"topic":"Git"}}\' | python notes_generator_agent.py',
            file=sys.stderr,
        )
        sys.exit(2)

    raw = raw.strip().lstrip("﻿")
    try:
        request_data = json.loads(raw)
    except json.JSONDecodeError as exc:
        result = {"status": "error", "action": "unknown", "message": f"Invalid JSON input: {exc}"}
        print(json.dumps(result, indent=2, ensure_ascii=False))
        sys.exit(1)

    result = execute(request_data)
    print(json.dumps(result, indent=2, ensure_ascii=False, default=str))
    sys.exit(0 if result.get("status") == "success" else 1)


if __name__ == "__main__":
    _cli_main()
