"""file_agent.py — File system and document access plugin (TOOL ONLY).

This module is a *pure tool*. It performs **no** analysis, summarization,
classification, reasoning, decision making, intent detection, or content
interpretation. It only:

    1. Receives a JSON-compatible ``dict`` request.
    2. Performs file system operations or reads documents.
    3. Extracts raw content and metadata.
    4. Returns a JSON-compatible ``dict`` response.

All intelligence belongs to the calling "Brain Agent". The single public
entry point is :func:`execute`.

Dependencies (install once)::

    pip install PyPDF2 python-docx openpyxl python-pptx pandas

CLI usage::

    python file_agent.py request.json
    python file_agent.py '{"action": "list_directory", "parameters": {"path": "."}}'
"""

from __future__ import annotations

import fnmatch
import json
import os
import sys
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Optional

__all__ = ["execute", "PLUGIN_INFO"]


# --------------------------------------------------------------------------- #
# Plugin metadata
# --------------------------------------------------------------------------- #
PLUGIN_INFO: dict[str, str] = {
    "name": "file_agent",
    "agent_name": "FileAgent",
    "version": "1.0.0",
    "type": "tool",
    "input_format": "json",
    "output_format": "json",
    "entrypoint": "execute",
    "description": "File system and document access plugin",
}


# --------------------------------------------------------------------------- #
# Supported extensions
# --------------------------------------------------------------------------- #
TEXT_EXTENSIONS: frozenset[str] = frozenset({
    ".txt", ".md", ".csv", ".json", ".py", ".js", ".ts",
    ".html", ".css", ".xml", ".yaml", ".yml",
})

DOCUMENT_EXTENSIONS: frozenset[str] = frozenset({
    ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx",
})

ALL_SUPPORTED_EXTENSIONS: frozenset[str] = TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS

DEFAULT_MAX_DEPTH: int = 5
DEFAULT_LIMIT: int = 20
MAX_TEXT_SIZE: int = 50 * 1024 * 1024  # 50 MB safety cap for text files
MAX_RECURSIVE_ENTRIES: int = 10_000


# --------------------------------------------------------------------------- #
# Exceptions
# --------------------------------------------------------------------------- #
class FileAgentError(Exception):
    """Raised for any handled error whose message is safe to return.

    Covers parameter validation, missing files, permission errors, and
    corrupt documents. These map to a structured ``status: error`` response
    instead of crashing the tool.
    """


# --------------------------------------------------------------------------- #
# Path helpers
# --------------------------------------------------------------------------- #
def _resolve_path(raw: str) -> Path:
    """Resolve and validate a user-supplied path string.

    Raises:
        FileAgentError: On empty input.
    """
    if not raw or not raw.strip():
        raise FileAgentError("Path parameter is empty")
    return Path(raw).resolve()


def _require_file(path: Path) -> None:
    """Assert that *path* exists and is a file."""
    if not path.exists():
        raise FileAgentError(f"File not found: {path}")
    if not path.is_file():
        raise FileAgentError(f"Path is not a file: {path}")


def _require_directory(path: Path) -> None:
    """Assert that *path* exists and is a directory."""
    if not path.exists():
        raise FileAgentError(f"Directory not found: {path}")
    if not path.is_dir():
        raise FileAgentError(f"Path is not a directory: {path}")


def _iso_timestamp(epoch: float) -> str:
    """Convert an epoch timestamp to an ISO-8601 UTC string."""
    return datetime.fromtimestamp(epoch, tz=timezone.utc).isoformat()


def _file_meta(path: Path) -> dict[str, Any]:
    """Return a standard metadata dict for a single file."""
    stat = path.stat()
    return {
        "name": path.name,
        "extension": path.suffix.lower(),
        "size_bytes": stat.st_size,
        "created_time": _iso_timestamp(stat.st_ctime),
        "modified_time": _iso_timestamp(stat.st_mtime),
        "absolute_path": str(path),
    }


# --------------------------------------------------------------------------- #
# Directory listing helpers
# --------------------------------------------------------------------------- #
def _list_entries(directory: Path) -> dict[str, list[dict[str, Any]]]:
    """List immediate children of *directory* split into folders and files."""
    folders: list[dict[str, Any]] = []
    files: list[dict[str, Any]] = []
    try:
        entries = sorted(directory.iterdir(), key=lambda p: p.name.lower())
    except PermissionError:
        raise FileAgentError(f"Permission denied: {directory}")

    for entry in entries:
        try:
            if entry.is_dir():
                folders.append({"name": entry.name, "path": str(entry)})
            elif entry.is_file():
                stat = entry.stat()
                files.append({
                    "name": entry.name,
                    "extension": entry.suffix.lower(),
                    "size_bytes": stat.st_size,
                    "modified_time": _iso_timestamp(stat.st_mtime),
                    "path": str(entry),
                })
        except (PermissionError, OSError):
            continue
    return {"folders": folders, "files": files}


def _walk_recursive(
    directory: Path, max_depth: int, _current_depth: int = 0
) -> dict[str, Any]:
    """Build a recursive directory tree up to *max_depth* levels."""
    node: dict[str, Any] = {
        "name": directory.name,
        "path": str(directory),
        "folders": [],
        "files": [],
    }
    if _current_depth >= max_depth:
        return node

    try:
        entries = sorted(directory.iterdir(), key=lambda p: p.name.lower())
    except PermissionError:
        node["error"] = "permission denied"
        return node

    for entry in entries:
        try:
            if entry.is_dir():
                child = _walk_recursive(entry, max_depth, _current_depth + 1)
                node["folders"].append(child)
            elif entry.is_file():
                stat = entry.stat()
                node["files"].append({
                    "name": entry.name,
                    "extension": entry.suffix.lower(),
                    "size_bytes": stat.st_size,
                    "modified_time": _iso_timestamp(stat.st_mtime),
                    "path": str(entry),
                })
        except (PermissionError, OSError):
            continue
    return node


# --------------------------------------------------------------------------- #
# Document readers
# --------------------------------------------------------------------------- #
def _read_text(path: Path) -> dict[str, Any]:
    """Read a text-based file and return its content."""
    stat = path.stat()
    if stat.st_size > MAX_TEXT_SIZE:
        raise FileAgentError(
            f"File too large ({stat.st_size:,} bytes, limit {MAX_TEXT_SIZE:,})"
        )
    try:
        content = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            content = path.read_text(encoding="latin-1")
        except Exception as exc:
            raise FileAgentError(f"Unable to decode file: {exc}")
    except PermissionError:
        raise FileAgentError(f"Permission denied: {path}")
    return {"content": content}


def _read_pdf(path: Path) -> dict[str, Any]:
    """Read a PDF and return per-page text."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        raise FileAgentError(
            "PyPDF2 is not installed. Run: pip install PyPDF2"
        )
    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        raise FileAgentError(f"Failed to open PDF: {exc}")

    pages: list[dict[str, Any]] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        pages.append({"page": i, "text": text})
    return {"page_count": len(pages), "pages": pages}


def _read_excel(path: Path) -> dict[str, Any]:
    """Read an Excel workbook and return per-sheet rows."""
    ext = path.suffix.lower()
    try:
        import pandas as pd
    except ImportError:
        raise FileAgentError(
            "pandas is not installed. Run: pip install pandas"
        )

    engine: Optional[str] = None
    if ext == ".xlsx":
        try:
            import openpyxl  # noqa: F401
            engine = "openpyxl"
        except ImportError:
            raise FileAgentError(
                "openpyxl is not installed. Run: pip install openpyxl"
            )
    elif ext == ".xls":
        try:
            import xlrd  # noqa: F401
            engine = "xlrd"
        except ImportError:
            raise FileAgentError(
                "xlrd is not installed. Run: pip install xlrd"
            )
    else:
        raise FileAgentError(f"Unsupported Excel extension: {ext}")

    try:
        sheets_dict = pd.read_excel(str(path), sheet_name=None, engine=engine)
    except Exception as exc:
        raise FileAgentError(f"Failed to read Excel file: {exc}")

    sheets: list[dict[str, Any]] = []
    for sheet_name, df in sheets_dict.items():
        df = df.fillna("")
        rows = df.values.tolist()
        columns = df.columns.tolist()
        sheets.append({
            "sheet_name": str(sheet_name),
            "columns": [str(c) for c in columns],
            "row_count": len(rows),
            "rows": rows,
        })
    return {"sheet_count": len(sheets), "sheets": sheets}


def _read_word(path: Path) -> dict[str, Any]:
    """Read a Word document and return its paragraphs."""
    ext = path.suffix.lower()
    if ext == ".doc":
        raise FileAgentError(
            "Legacy .doc format is not supported. Convert to .docx first."
        )
    try:
        from docx import Document
    except ImportError:
        raise FileAgentError(
            "python-docx is not installed. Run: pip install python-docx"
        )
    try:
        doc = Document(str(path))
    except Exception as exc:
        raise FileAgentError(f"Failed to open Word document: {exc}")

    paragraphs: list[str] = [p.text for p in doc.paragraphs]
    return {"paragraph_count": len(paragraphs), "paragraphs": paragraphs}


def _read_powerpoint(path: Path) -> dict[str, Any]:
    """Read a PowerPoint presentation and return per-slide text."""
    ext = path.suffix.lower()
    if ext == ".ppt":
        raise FileAgentError(
            "Legacy .ppt format is not supported. Convert to .pptx first."
        )
    try:
        from pptx import Presentation
    except ImportError:
        raise FileAgentError(
            "python-pptx is not installed. Run: pip install python-pptx"
        )
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise FileAgentError(f"Failed to open PowerPoint file: {exc}")

    slides: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        slides.append({"slide_number": i, "text": "\n".join(texts)})
    return {"slide_count": len(slides), "slides": slides}


def _read_any_file(path: Path) -> dict[str, Any]:
    """Read a file based on its extension, returning structured content."""
    ext = path.suffix.lower()
    if ext in TEXT_EXTENSIONS:
        return _read_text(path)
    if ext == ".pdf":
        return _read_pdf(path)
    if ext in (".xlsx", ".xls"):
        return _read_excel(path)
    if ext in (".docx", ".doc"):
        return _read_word(path)
    if ext in (".pptx", ".ppt"):
        return _read_powerpoint(path)
    raise FileAgentError(f"Unsupported file type: {ext}")


# --------------------------------------------------------------------------- #
# Action handlers
# --------------------------------------------------------------------------- #
def _action_list_directory(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_directory(path)
    return _list_entries(path)


def _action_list_directory_recursive(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_directory(path)
    max_depth = int(params.get("max_depth", DEFAULT_MAX_DEPTH))
    if max_depth < 1:
        raise FileAgentError("max_depth must be at least 1")
    return _walk_recursive(path, max_depth)


def _action_search_files(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_directory(path)
    pattern = params.get("pattern", "")
    if not pattern or not pattern.strip():
        raise FileAgentError("Missing pattern parameter")

    matches: list[dict[str, Any]] = []
    count = 0
    for root, _dirs, files in os.walk(str(path)):
        for name in files:
            if fnmatch.fnmatch(name.lower(), pattern.lower()):
                fp = Path(root) / name
                try:
                    stat = fp.stat()
                    matches.append({
                        "name": name,
                        "extension": fp.suffix.lower(),
                        "size_bytes": stat.st_size,
                        "modified_time": _iso_timestamp(stat.st_mtime),
                        "path": str(fp),
                    })
                except (PermissionError, OSError):
                    continue
                count += 1
                if count >= MAX_RECURSIVE_ENTRIES:
                    break
        if count >= MAX_RECURSIVE_ENTRIES:
            break
    return {"count": len(matches), "matches": matches}


def _action_get_file_info(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_file(path)
    return _file_meta(path)


def _action_read_text_file(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_file(path)
    ext = path.suffix.lower()
    if ext not in TEXT_EXTENSIONS:
        raise FileAgentError(
            f"Unsupported text file type: {ext}. "
            f"Supported: {', '.join(sorted(TEXT_EXTENSIONS))}"
        )
    return _read_text(path)


def _action_read_pdf(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_file(path)
    if path.suffix.lower() != ".pdf":
        raise FileAgentError(f"Not a PDF file: {path.suffix}")
    return _read_pdf(path)


def _action_read_excel(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_file(path)
    ext = path.suffix.lower()
    if ext not in (".xlsx", ".xls"):
        raise FileAgentError(f"Not an Excel file: {ext}")
    return _read_excel(path)


def _action_read_word(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_file(path)
    ext = path.suffix.lower()
    if ext not in (".docx", ".doc"):
        raise FileAgentError(f"Not a Word document: {ext}")
    return _read_word(path)


def _action_read_powerpoint(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_file(path)
    ext = path.suffix.lower()
    if ext not in (".pptx", ".ppt"):
        raise FileAgentError(f"Not a PowerPoint file: {ext}")
    return _read_powerpoint(path)


def _action_read_multiple_files(params: dict[str, Any]) -> dict[str, Any]:
    paths = params.get("paths")
    if not paths or not isinstance(paths, list):
        raise FileAgentError("Missing or invalid paths parameter (expected a list)")
    if not all(isinstance(p, str) and p.strip() for p in paths):
        raise FileAgentError("All items in paths must be non-empty strings")

    results: list[dict[str, Any]] = []
    for raw in paths:
        try:
            path = _resolve_path(raw)
            _require_file(path)
            content = _read_any_file(path)
            results.append({
                "status": "success",
                "path": str(path),
                "name": path.name,
                "extension": path.suffix.lower(),
                "data": content,
            })
        except (FileAgentError, Exception) as exc:
            results.append({
                "status": "error",
                "path": raw,
                "message": str(exc),
            })
    return {"count": len(results), "results": results}


def _action_get_folder_tree(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_directory(path)
    max_depth = int(params.get("max_depth", DEFAULT_MAX_DEPTH))
    if max_depth < 1:
        raise FileAgentError("max_depth must be at least 1")
    return _walk_recursive(path, max_depth)


def _action_file_exists(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    exists = path.is_file()
    result: dict[str, Any] = {"exists": exists, "path": str(path)}
    if exists:
        result["name"] = path.name
        result["extension"] = path.suffix.lower()
        result["size_bytes"] = path.stat().st_size
    return result


def _action_directory_exists(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    exists = path.is_dir()
    return {"exists": exists, "path": str(path)}


def _action_get_supported_files(params: dict[str, Any]) -> dict[str, Any]:
    return {
        "text_extensions": sorted(TEXT_EXTENSIONS),
        "document_extensions": sorted(DOCUMENT_EXTENSIONS),
        "all_extensions": sorted(ALL_SUPPORTED_EXTENSIONS),
    }


def _action_get_recent_files(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_directory(path)
    limit = int(params.get("limit", DEFAULT_LIMIT))
    if limit < 1:
        raise FileAgentError("limit must be at least 1")

    all_files: list[tuple[float, Path]] = []
    count = 0
    for root, _dirs, files in os.walk(str(path)):
        for name in files:
            fp = Path(root) / name
            try:
                mtime = fp.stat().st_mtime
                all_files.append((mtime, fp))
            except (PermissionError, OSError):
                continue
            count += 1
            if count >= MAX_RECURSIVE_ENTRIES:
                break
        if count >= MAX_RECURSIVE_ENTRIES:
            break

    all_files.sort(key=lambda t: t[0], reverse=True)
    recent = all_files[:limit]

    results: list[dict[str, Any]] = []
    for mtime, fp in recent:
        try:
            stat = fp.stat()
            results.append({
                "name": fp.name,
                "extension": fp.suffix.lower(),
                "size_bytes": stat.st_size,
                "modified_time": _iso_timestamp(stat.st_mtime),
                "path": str(fp),
            })
        except (PermissionError, OSError):
            continue
    return {"count": len(results), "files": results}


def _action_read_folder_contents(params: dict[str, Any]) -> dict[str, Any]:
    path = _resolve_path(params.get("path", ""))
    _require_directory(path)
    limit = int(params.get("limit", 50))

    results: list[dict[str, Any]] = []
    try:
        entries = sorted(path.iterdir(), key=lambda p: p.name.lower())
    except PermissionError:
        raise FileAgentError(f"Permission denied: {path}")

    for entry in entries:
        if not entry.is_file():
            continue
        ext = entry.suffix.lower()
        if ext not in ALL_SUPPORTED_EXTENSIONS:
            continue
        try:
            content = _read_any_file(entry)
            results.append({
                "status": "success",
                "path": str(entry),
                "name": entry.name,
                "extension": ext,
                "data": content,
            })
        except (FileAgentError, Exception) as exc:
            results.append({
                "status": "error",
                "path": str(entry),
                "name": entry.name,
                "message": str(exc),
            })
        if len(results) >= limit:
            break
    return {"count": len(results), "results": results}


# --------------------------------------------------------------------------- #
# Action registry and parameter validation
# --------------------------------------------------------------------------- #
_REQUIRED_PARAMS: dict[str, list[str]] = {
    "list_directory": ["path"],
    "list_directory_recursive": ["path"],
    "search_files": ["path", "pattern"],
    "get_file_info": ["path"],
    "read_text_file": ["path"],
    "read_pdf": ["path"],
    "read_excel": ["path"],
    "read_word": ["path"],
    "read_powerpoint": ["path"],
    "read_multiple_files": ["paths"],
    "get_folder_tree": ["path"],
    "file_exists": ["path"],
    "directory_exists": ["path"],
    "get_supported_files": [],
    "get_recent_files": ["path"],
    "read_folder_contents": ["path"],
}

_ACTIONS: dict[str, Callable[[dict[str, Any]], dict[str, Any]]] = {
    "list_directory": _action_list_directory,
    "list_directory_recursive": _action_list_directory_recursive,
    "search_files": _action_search_files,
    "get_file_info": _action_get_file_info,
    "read_text_file": _action_read_text_file,
    "read_pdf": _action_read_pdf,
    "read_excel": _action_read_excel,
    "read_word": _action_read_word,
    "read_powerpoint": _action_read_powerpoint,
    "read_multiple_files": _action_read_multiple_files,
    "get_folder_tree": _action_get_folder_tree,
    "file_exists": _action_file_exists,
    "directory_exists": _action_directory_exists,
    "get_supported_files": _action_get_supported_files,
    "get_recent_files": _action_get_recent_files,
    "read_folder_contents": _action_read_folder_contents,
}


# --------------------------------------------------------------------------- #
# Public entry point
# --------------------------------------------------------------------------- #
def execute(request_json: dict[str, Any]) -> dict[str, Any]:
    """Execute a file system or document-reading action.

    Args:
        request_json: A dict with ``"action"`` (str) and ``"parameters"`` (dict).

    Returns:
        A JSON-compatible dict with ``"status"`` (``"success"`` | ``"error"``),
        ``"action"``, and either ``"data"`` or ``"message"``.
    """
    if isinstance(request_json, str):
        request_json = request_json.lstrip("﻿")
        try:
            request_json = json.loads(request_json)
        except json.JSONDecodeError as exc:
            return {"status": "error", "action": "unknown", "message": f"Invalid JSON: {exc}"}

    if not isinstance(request_json, dict):
        return {"status": "error", "action": "unknown", "message": "Request must be a JSON object"}

    action = request_json.get("action", "")
    if not action or not isinstance(action, str):
        return {"status": "error", "action": "unknown", "message": "Missing or invalid action field"}

    if action not in _ACTIONS:
        return {
            "status": "error",
            "action": action,
            "message": f"Unknown action: {action}. Available: {', '.join(sorted(_ACTIONS))}",
        }

    parameters = request_json.get("parameters", {})
    if not isinstance(parameters, dict):
        return {"status": "error", "action": action, "message": "parameters must be a JSON object"}

    required = _REQUIRED_PARAMS.get(action, [])
    for key in required:
        value = parameters.get(key)
        if value is None or (isinstance(value, str) and not value.strip()):
            return {"status": "error", "action": action, "message": f"Missing required parameter: {key}"}

    try:
        data = _ACTIONS[action](parameters)
        return {"status": "success", "action": action, "data": data}
    except FileAgentError as exc:
        return {"status": "error", "action": action, "message": str(exc)}
    except PermissionError as exc:
        return {"status": "error", "action": action, "message": f"Permission denied: {exc}"}
    except Exception as exc:
        return {"status": "error", "action": action, "message": f"Unexpected error: {exc}"}


# --------------------------------------------------------------------------- #
# CLI entry point
# --------------------------------------------------------------------------- #
def _cli_main() -> None:
    """Handle command-line invocation: file arg, inline JSON arg, or piped stdin."""
    raw: Optional[str] = None

    if len(sys.argv) > 1:
        arg = sys.argv[1]
        if os.path.isfile(arg):
            with open(arg, "rb") as fh:
                raw = fh.read().decode("utf-8-sig")
        else:
            raw = arg
    elif not sys.stdin.isatty():
        raw = sys.stdin.buffer.read().decode("utf-8-sig")
    else:
        print(
            "Usage:\n"
            '  python file_agent.py request.json\n'
            '  python file_agent.py \'{"action":"list_directory","parameters":{"path":"."}}\'\n'
            '  echo \'{"action":"file_exists","parameters":{"path":"test.txt"}}\' | python file_agent.py',
            file=sys.stderr,
        )
        sys.exit(2)

    raw = raw.strip().lstrip("﻿")
    try:
        request_data = json.loads(raw)
    except json.JSONDecodeError as exc:
        result = {"status": "error", "action": "unknown", "message": f"Invalid JSON input: {exc}"}
        print(json.dumps(result, indent=2, ensure_ascii=False))
        sys.exit(1)

    result = execute(request_data)
    print(json.dumps(result, indent=2, ensure_ascii=False, default=str))
    sys.exit(0 if result.get("status") == "success" else 1)


if __name__ == "__main__":
    _cli_main()
